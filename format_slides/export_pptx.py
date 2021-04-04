from pptx import Presentation
from pptx.util import Inches
import time
from selenium import webdriver
from PIL import Image
import io
import requests
import yfinance as yf
import requests
import wikipedia

# Get user input for ticker as variable 'ticker'
ticker = input('Enter a ticker symbol: ').upper()
# Initialise webdriver for selenium
driver = webdriver.Chrome('/usr/local/bin/chromedriver') 

# -------------- You can change these: -------------- #
# Input and output pptx file names (You can change these but make sure the template file is in the same folder as this python script)
template_name = 'Template.pptx'
output_name = ticker + ' Black Boar Capital Pitch.pptx'



# -------------- Functions: -------------- #
def get_full_name(ticker_in):
    url = "http://d.yimg.com/autoc.finance.yahoo.com/autoc?query={}&region=1&lang=en".format(ticker_in)

    result = requests.get(url).json()

    for x in result['ResultSet']['Result']:
        if x['symbol'] == ticker_in:
            return x['name']

def replace_text(find_in, replace_in):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                if(shape.text.find(find_in))!=-1:
                    print('Replaced!')
                    text_frame = shape.text_frame
                    cur_text = text_frame.paragraphs[0].runs[0].text
                    new_text = cur_text.replace(find_in, replace_in) # Automate finding full company name
                    text_frame.paragraphs[0].runs[0].text = new_text


def replace_logo():
	for slide in prs.slides:
	    for shape in slide.shapes:
	        if shape.has_text_frame:
	            if(shape.text.find('Company_logo'))!=-1:
	                pic = slide.shapes.add_picture(r'logo.jpg', shape.left, shape.top, height=Inches(1.5))


def fetch_image_urls(query:str, max_links_to_fetch:int=1, wd=driver, sleep_between_interactions:int=1):
    def scroll_to_end(wd):
        wd.execute_script("window.scrollTo(0, document.body.scrollHeight);")
        time.sleep(sleep_between_interactions)    
    
    # build the google query
    search_url = "https://www.google.com/search?safe=off&site=&tbm=isch&source=hp&q={q}&oq={q}&gs_l=img"

    # load the page
    wd.get(search_url.format(q=query))

    image_urls = set()
    image_count = 0
    results_start = 0
    while image_count < max_links_to_fetch:
        scroll_to_end(wd)

        # get all image thumbnail results
        thumbnail_results = wd.find_elements_by_css_selector("img.Q4LuWd")
        number_results = len(thumbnail_results)
        
        print(f"Found: {number_results} search results. Extracting links from {results_start}:{number_results}")
        
        for img in thumbnail_results[results_start:number_results]:
            # try to click every thumbnail such that we can get the real image behind it
            try:
                img.click()
                time.sleep(sleep_between_interactions)
            except Exception:
                continue

            # extract image urls    
            actual_images = wd.find_elements_by_css_selector('img.n3VNCb')
            for actual_image in actual_images:
                if actual_image.get_attribute('src') and 'http' in actual_image.get_attribute('src'):
                    image_urls.add(actual_image.get_attribute('src'))

            image_count = len(image_urls)

            if len(image_urls) >= max_links_to_fetch:
                print(f"Found: {len(image_urls)} image links, done!")
                break
        else:
            print("Found:", len(image_urls), "image links, looking for more ...")
            time.sleep(30)
            return
            load_more_button = wd.find_element_by_css_selector(".mye4qd")
            if load_more_button:
                wd.execute_script("document.querySelector('.mye4qd').click();")

        # move the result startpoint further down
        results_start = len(thumbnail_results)

    return image_urls

def persist_image(url:str):
    try:
        image_content = requests.get(url).content

    except Exception as e:
        print(f"ERROR - Could not download {url} - {e}")

    try:
        image_file = io.BytesIO(image_content)
        image = Image.open(image_file).convert('RGB')
        file_path = 'logo.jpg'
        with open(file_path, 'wb') as f:
            image.save(f, "JPEG", quality=85)
        print(f"SUCCESS - saved {url} - as {file_path}")
    except Exception as e:
        print(f"ERROR - Could not save {url} - {e}")

    driver.quit()


def save_changes():
	prs.save(output_name)



# --------------- Main --------------- #
# Initialise presentation
prs = Presentation(template_name)
# Get ticker data from Yahoo! Finance for sector information (only sector information for now)
tickerdata = yf.Ticker(ticker)
# Get full name of company from ticker symbol
full_company_name = get_full_name(ticker)

# Format slides to include the company name, sector and description
replace_text('Company_name', full_company_name + ' (' + ticker + ')')
replace_text('Company_sector', tickerdata.info['sector'])
replace_text('Company_summary', (wikipedia.summary(full_company_name +' (company)')))

# Download logo
for i in fetch_image_urls(full_company_name + ' logo'):
	persist_image(i)

# Include logo
replace_logo()

# Save changes to powerpoint as new file
save_changes()

# Very Nice!
print('Great success!')





